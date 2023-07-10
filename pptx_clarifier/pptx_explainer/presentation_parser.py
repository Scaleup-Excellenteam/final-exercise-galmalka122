import argparse
import os
import re
from typing import Optional

from pptx import Presentation, exc

from pptx_clarifier.pptx_explainer import explainer_logger as logger


def open_presentation(path: str) -> Presentation:
    """Opens a presentation file asynchronously

    Args:
        path (str): path to the presentation file

    Returns:
        Presentation: presentation object

    Raises:
        FileNotFoundError: if the path does not exist
        ValueError: if the path is not a pptx file
    """
    validate_path(path)
    presentation = Presentation(path)
    return presentation


def parse_slide(slide) -> Optional[dict]:
    """Gets the title and the text from all shapes in a slide

    Args:
        slide (slide): slide to be parsed

    Returns:
        dict: dictionary with the title and the text of the slide

    Raises:
        Exception: if an error occurs while parsing the slide
    """
    try:
        title = None
        if hasattr(slide.shapes, 'title'):
            title = remove_special_characters(slide.shapes.title.text)
        slide_text = [remove_special_characters(shape.text)
                      for shape in slide.shapes
                      if shape.has_text_frame and shape.text]
        if len(slide_text) == 0 and not title:
            return None
        if title in slide_text:
            slide_text.remove(title)
        slide_data = {"title": title, "text": slide_text}
        return slide_data
    except exc.PythonPptxError as pptx_exc:
        logger.error(f"An error occurred while parsing slide: {pptx_exc}")
        return None


def remove_special_characters(text):
    """Removes special characters from text and replaces them with spaces
    if the character is a newline, it is replaced with a literal \n

    Args:
        text (str): text to be cleaned

    Returns:
        str: cleaned text
    """
    text = re.sub('[^A-Za-z0-9]+', ' ', text)
    text = re.sub(r'\n', r'\\n', text)
    return text


def validate_path(path) -> None:
    """Validates the path is a real path and a pptx file

    Args:
        path (str): path to be validated

    Returns:
        str: error message if the path is not valid
        function: presentation_parser function if the path is valid

    Raises:
        FileNotFoundError: if the path does not exist
        ValueError: if the path is not a pptx file
    """
    # check if path is local or absolute path and convert to absolute path
    if not os.path.isabs(path):
        path = os.path.abspath(path)
    if not os.path.exists(path):
        logger.error('User entered an invalid path')
        raise FileNotFoundError(f'{path} does not exist')
    if not path.endswith('.pptx'):
        raise ValueError(f'{path} is not a pptx file')


def main():
    # Create an argument parser
    parser = argparse.ArgumentParser(description='Script to parse pptx file to text.')

    # Add an argument for the file path
    parser.add_argument('file_path', nargs=1, type=str, help='Path to the pptx file')

    # Parse the command-line arguments
    args = parser.parse_args()

    # Extract the file path from the arguments
    file_path = args.file_path

    # Check if the file path argument is missing
    if not file_path:
        parser.error('Please provide a path to the pptx file.')
    try:
        # Call the function to process the file
        prs = open_presentation(file_path)
        data = [parse_slide(sl) for sl in prs.slides]
        print(data)
    except Exception as e:
        print(e)


if __name__ == '__main__':
    main()
