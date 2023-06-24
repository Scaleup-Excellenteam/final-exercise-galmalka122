import os
import re
import asyncio

from pptx import Presentation

from .logger import logger


async def open_presentation(path: str):
    """Opens a presentation file asynchronously

    Args:
        path (str): path to the presentation file

    Returns:
        Presentation: presentation object

    Raises:
        FileNotFoundError: if the path does not exist
        ValueError: if the path is not a pptx file
    """
    await validate_path(path)
    prs = await asyncio.get_running_loop().run_in_executor(None, Presentation, path)
    return prs


def parse_slide(slide):
    """Gets the title and the text from all shapes in a slide

    Args:
        slide (slide): slide to be parsed

    Returns:
        dict: dictionary with the title and the text of the slide

    Raises:
        Exception: if an error occurs while parsing the slide
    """
    try:
        title = remove_special_characters(slide.shapes.title.text)
        slide_text = [remove_special_characters(shape.text)
                      for shape in slide.shapes
                      if shape.has_text_frame
                      and shape.text != ""]
        if title in slide_text:
            slide_text.remove(title)
        data = {"title": title, "text": slide_text}
        return data
    except Exception as e:
        logger.error(f"Error in parse_presentation: {e} occurred")
        raise RuntimeError("An error occurred while parsing the presentation.")


def remove_special_characters(text):
    """Removes special characters from text and replaces them with spaces
    if the character is a newline, it is replaced with a literal \n

    Args:
        text (str): text to be cleaned

    Returns:
        str: cleaned text
    """
    text = re.sub('[^A-Za-z0-9]+', ' ', text)
    text = re.sub(r'\n', r'\\n', text)
    return text


async def validate_path(path) -> None:
    """Validates the path is a real path and a pptx file

    Args:
        path (str): path to be validated

    Returns:
        str: error message if the path is not valid
        function: presentation_parser function if the path is valid

    Raises:
        FileNotFoundError: if the path does not exist
        ValueError: if the path is not a pptx file
    """
    # check if path is local or absolute path and convert to absolute path
    if not os.path.isabs(path):
        path = os.path.abspath(path)
    if not os.path.exists(path):
        logger.error('User entered an invalid path')
        raise FileNotFoundError(f'{path} does not exist')
